import os
import re
import datetime
import uuid
import pypdf
import docx
import openpyxl
import zipfile
from collections import Counter
from typing import List, Dict, Any, Tuple, Optional
from app.schemas.normalized_block import NormalizedBlock
from app.core.logging import logger

try:
    from markitdown import MarkItDown
    markitdown_available = True
    markitdown_instance = MarkItDown()
except ImportError:
    markitdown_available = False
    markitdown_instance = None
    logger.warning("markitdown package not found, falling back to direct parsing.")

class ParserService:
    @staticmethod
    def _save_image_with_hash(
        img_bytes: bytes,
        doc_basename: str,
        ext: str,
        base_dir: str,
        doc_image_cache: Dict[str, Dict[str, Any]],
        alt_text: str = ""
    ) -> Dict[str, Any]:
        """
        Saves image bytes using MD5 hash deduplication.
        If an identical image byte stream was already saved for this document/session:
        - Does NOT write a new duplicate file to disk.
        - Reuses the existing file paths and metadata.
        Returns dict with 'rel_path', 'full_path', 'alt', 'hash', 'is_duplicate', 'count'.
        """
        import hashlib
        img_hash = hashlib.md5(img_bytes).hexdigest()
        ext = ext.lower().lstrip(".")
        if ext not in ["png", "jpeg", "jpg", "gif", "svg", "webp"]:
            ext = "png"

        if img_hash in doc_image_cache:
            cached = doc_image_cache[img_hash]
            cached["count"] += 1
            return {
                "rel_path": cached["rel_path"],
                "full_path": cached["full_path"],
                "alt": alt_text or cached["alt"],
                "hash": img_hash,
                "is_duplicate": True,
                "count": cached["count"]
            }

        img_filename = f"{doc_basename}_img_{img_hash[:10]}.{ext}"
        img_full_path = os.path.join(base_dir, img_filename)
        rel_path = f"storage/extracted_images/{img_filename}"

        with open(img_full_path, "wb") as f_img:
            f_img.write(img_bytes)

        res = {
            "rel_path": rel_path,
            "full_path": img_full_path,
            "alt": alt_text,
            "hash": img_hash,
            "is_duplicate": False,
            "count": 1
        }
        doc_image_cache[img_hash] = res
        return res

    @staticmethod
    def parse_markdown_content(
        md_text: str, 
        source_type: str, 
        page_start: Optional[int] = None, 
        page_end: Optional[int] = None,
        initial_order: int = 0
    ) -> List[NormalizedBlock]:
        """
        Parses standardized Markdown text into structured NormalizedBlock list.
        Tracks heading hierarchy (heading_path), code blocks, tables, and lists.
        """
        blocks: List[NormalizedBlock] = []
        lines = md_text.split("\n")
        
        current_heading_path: List[str] = []
        source_order = initial_order
        
        buf_type = "paragraph"
        buf_lines: List[str] = []
        in_code_block = False
        code_block_lines: List[str] = []

        def flush_buffer():
            nonlocal source_order, buf_lines, buf_type
            if not buf_lines:
                return
            
            content = "\n".join(buf_lines).strip()
            if content:
                source_order += 1
                blocks.append(NormalizedBlock(
                    block_id=str(uuid.uuid4()),
                    source_type=source_type,
                    block_type=buf_type,
                    content=content,
                    heading_path=list(current_heading_path),
                    page_start=page_start,
                    page_end=page_end,
                    source_order=source_order
                ))
            buf_lines = []
            buf_type = "paragraph"

        for line in lines:
            trimmed = line.strip()

            # Filter out decorative horizontal rules (e.g. ---, ***, ___, --------...)
            if re.match(r"^[\s\-\*_]{3,}$", trimmed):
                flush_buffer()
                continue

            # Handle Code Blocks ```
            if trimmed.startswith("```"):
                if in_code_block:
                    code_block_lines.append(trimmed)
                    source_order += 1
                    blocks.append(NormalizedBlock(
                        block_id=str(uuid.uuid4()),
                        source_type=source_type,
                        block_type="code",
                        content="\n".join(code_block_lines),
                        heading_path=list(current_heading_path),
                        page_start=page_start,
                        page_end=page_end,
                        source_order=source_order
                    ))
                    code_block_lines = []
                    in_code_block = False
                else:
                    flush_buffer()
                    in_code_block = True
                    code_block_lines = [trimmed]
                continue

            if in_code_block:
                code_block_lines.append(line)
                continue

            # Handle Headings (# Heading)
            heading_match = re.match(r"^(#{1,6})\s+(.*)$", trimmed)
            if heading_match:
                flush_buffer()
                level = len(heading_match.group(1))
                h_text = heading_match.group(2).strip()

                if level == 1:
                    current_heading_path = [h_text]
                elif level == 2:
                    current_heading_path = [current_heading_path[0], h_text] if len(current_heading_path) >= 1 else [h_text]
                else:
                    current_heading_path = current_heading_path[:level-1] + [h_text]

                source_order += 1
                blocks.append(NormalizedBlock(
                    block_id=str(uuid.uuid4()),
                    source_type=source_type,
                    block_type="heading",
                    content=h_text,
                    heading_path=list(current_heading_path),
                    page_start=page_start,
                    page_end=page_end,
                    source_order=source_order
                ))
                continue

            # Handle Markdown Tables (| col | col |)
            if trimmed.startswith("|") and trimmed.endswith("|"):
                if buf_type != "table":
                    flush_buffer()
                    buf_type = "table"
                buf_lines.append(trimmed)
                continue
            elif buf_type == "table":
                flush_buffer()

            # Handle Lists (- item, * item, 1. item)
            if re.match(r"^(\*|-|\+|\d+\.)\s+", trimmed):
                if buf_type != "list":
                    flush_buffer()
                    buf_type = "list"
                buf_lines.append(trimmed)
                continue
            elif buf_type == "list" and trimmed == "":
                flush_buffer()

            # Empty lines flush buffer
            if not trimmed:
                flush_buffer()
                continue

            # Paragraph lines
            if buf_type not in ["paragraph", "list"]:
                flush_buffer()
                buf_type = "paragraph"
            buf_lines.append(line)

        flush_buffer()
        return blocks

    @classmethod
    def parse_pdf(cls, file_path: str) -> List[NormalizedBlock]:
        logger.info(f"Parsing PDF file with Layout Font-Size Parser: {file_path}")
        try:
            import pdfplumber
            from collections import Counter

            font_sizes = []
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    words = page.extract_words(extra_attrs=['fontname', 'size'])
                    for w in words:
                        font_sizes.append(round(w['size'], 1))

                if font_sizes:
                    base_size = Counter(font_sizes).most_common(1)[0][0]
                    md_lines = []

                    for page_idx, page in enumerate(pdf.pages):
                        page_num = page_idx + 1
                        words = page.extract_words(extra_attrs=['fontname', 'size'])
                        lines_dict = {}
                        for w in words:
                            top = round(w['top'], 1)
                            matching_top = None
                            for k in lines_dict:
                                if abs(k - top) <= 3:
                                    matching_top = k
                                    break
                            if matching_top is None:
                                matching_top = top
                                lines_dict[matching_top] = []
                            lines_dict[matching_top].append(w)

                        for top in sorted(lines_dict.keys()):
                            line_words = sorted(lines_dict[top], key=lambda x: x['x0'])
                            line_text = ' '.join([w['text'] for w in line_words]).strip()
                            if not line_text:
                                continue

                            avg_size = sum(w['size'] for w in line_words) / len(line_words)
                            is_bold = any('bold' in w['fontname'].lower() or 'heavy' in w['fontname'].lower() for w in line_words)

                            # Multi-Criteria Heading Promotion Engine (Font Size, ALL CAPS, Structural Heading Patterns)
                            heading_pattern = re.compile(
                                r"^(vụ việc:|mục\s+\d+|điều\s+\d+|chương\s+\d+|phần\s+\d+|báo cáo:|quy định:|quyết định:|hướng dẫn:|\d+\.|\d+\.\d+|[A-Z]\.|\b[I|V|X]+\.)\b",
                                re.IGNORECASE
                            )
                            is_all_caps = line_text.isupper() and 3 < len(line_text) < 120
                            is_title_pattern = bool(heading_pattern.search(line_text)) and len(line_text) < 120

                            if avg_size >= base_size * 1.35 or is_all_caps:
                                md_lines.append(f"# {line_text}")
                            elif avg_size >= base_size * 1.15 or is_bold or is_title_pattern:
                                md_lines.append(f"## {line_text}")
                            else:
                                md_lines.append(line_text)

                    full_md_text = "\n\n".join(md_lines)
                    return cls.parse_markdown_content(md_text=full_md_text, source_type="pdf")

        except Exception as err:
            logger.warning(f"pdfplumber layout parsing failed for {file_path}, falling back to pypdf: {str(err)}")

        # Fallback to pypdf
        blocks: List[NormalizedBlock] = []
        try:
            reader = pypdf.PdfReader(file_path)
            order_counter = 0

            for page_idx, page in enumerate(reader.pages):
                page_num = page_idx + 1
                text = page.extract_text() or ""
                if not text.strip():
                    continue

                page_blocks = cls.parse_markdown_content(
                    md_text=text,
                    source_type="pdf",
                    page_start=page_num,
                    page_end=page_num,
                    initial_order=order_counter
                )
                order_counter += len(page_blocks)
                blocks.extend(page_blocks)

        except Exception as e:
            logger.error(f"Error parsing PDF {file_path}: {str(e)}")
            raise e

        return blocks

    @classmethod
    def parse_docx(cls, file_path: str) -> List[NormalizedBlock]:
        logger.info(f"Parsing DOCX file with Table & Image extraction: {file_path}")
        blocks: List[NormalizedBlock] = []
        try:
            doc = docx.Document(file_path)
            source_order = 0
            current_heading_path = ["Tổng quan"]

            base_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "storage", "extracted_images"))
            os.makedirs(base_dir, exist_ok=True)
            doc_basename = os.path.splitext(os.path.basename(file_path))[0]
            doc_image_cache: Dict[str, Dict[str, Any]] = {}
            image_counter = 0

            # Iterate over elements in body sequentially (paragraphs and tables)
            for element in doc.element.body:
            # Handle Paragraph
                if element.tag.endswith('p'):
                    para = docx.text.paragraph.Paragraph(element, doc)
                    text = para.text.strip()
                    
                    # 1. Check for embedded images in this paragraph
                    para_img_tags = []
                    blips = element.xpath('.//a:blip')
                    for blip in blips:
                        embed_id = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                        if embed_id and embed_id in doc.part.rels:
                            rel_part = doc.part.rels[embed_id].target_part
                            if "image" in rel_part.content_type:
                                image_counter += 1
                                ext = rel_part.content_type.split("/")[-1]
                                img_info = cls._save_image_with_hash(
                                    img_bytes=rel_part.blob,
                                    doc_basename=doc_basename,
                                    ext=ext,
                                    base_dir=base_dir,
                                    doc_image_cache=doc_image_cache,
                                    alt_text=f"Embedded Image {image_counter} in {doc_basename}"
                                )
                                img_tag = f"<img src='{img_info['rel_path']}' alt='{img_info['alt']}' />"
                                para_img_tags.append(img_tag)

                    # Insert inline image tags into text paragraph for contextual linkage
                    if para_img_tags:
                        img_combined_tag = "\n\n" + "\n\n".join(para_img_tags)
                        text = (text + img_combined_tag).strip() if text else img_combined_tag.strip()

                    if not text:
                        continue

                    style_name = para.style.name.lower() if para.style else ""
                    if "heading" in style_name or text.startswith("#"):
                        heading_text = re.sub(r"^\#+\s*", "", text).strip()
                        current_heading_path = [heading_text]
                        source_order += 1
                        blocks.append(NormalizedBlock(
                            block_id=str(uuid.uuid4()),
                            source_type="docx",
                            block_type="heading",
                            content=heading_text,
                            heading_path=list(current_heading_path),
                            source_order=source_order
                        ))
                    else:
                        source_order += 1
                        blocks.append(NormalizedBlock(
                            block_id=str(uuid.uuid4()),
                            source_type="docx",
                            block_type="paragraph",
                            content=text,
                            heading_path=list(current_heading_path),
                            source_order=source_order
                        ))

                # Handle Table (Markdown Pipe Table format for token optimization)
                elif element.tag.endswith('tbl'):
                    table = docx.table.Table(element, doc)
                    table_rows = []
                    for row in table.rows:
                        row_vals = [cell.text.strip().replace('\n', ' ').replace('|', '\\|') for cell in row.cells]
                        if any(row_vals):
                            table_rows.append(row_vals)

                    if table_rows:
                        headers = table_rows[0]
                        data_rows = table_rows[1:]
                        
                        # Build Markdown Pipe header & divider
                        header_line = "| " + " | ".join(headers) + " |\n"
                        divider_line = "| " + " | ".join(["---"] * len(headers)) + " |\n"
                        header_md = header_line + divider_line

                        # Batch data rows into sub-tables (~15-20 rows or ~1000 chars per sub-table block)
                        curr_rows_md = []
                        curr_char_len = len(header_md)

                        for r_cols in data_rows:
                            row_line = "| " + " | ".join(r_cols) + " |\n"
                            row_len = len(row_line)

                            if curr_char_len + row_len >= 1000 or len(curr_rows_md) >= 20:
                                source_order += 1
                                sub_table_md = header_md + "".join(curr_rows_md)
                                blocks.append(NormalizedBlock(
                                    block_id=str(uuid.uuid4()),
                                    source_type="docx",
                                    block_type="table",
                                    content=sub_table_md.strip(),
                                    requires_llm_summary=True,
                                    heading_path=list(current_heading_path),
                                    source_order=source_order
                                ))
                                curr_rows_md = [row_line]
                                curr_char_len = len(header_md) + row_len
                            else:
                                curr_rows_md.append(row_line)
                                curr_char_len += row_len

                        # Flush remaining table rows batch
                        if curr_rows_md or not data_rows:
                            source_order += 1
                            sub_table_md = header_md + "".join(curr_rows_md)
                            blocks.append(NormalizedBlock(
                                block_id=str(uuid.uuid4()),
                                source_type="docx",
                                block_type="table",
                                content=sub_table_md.strip(),
                                requires_llm_summary=True,
                                heading_path=list(current_heading_path),
                                source_order=source_order
                            ))

        except Exception as e:
            logger.error(f"Error parsing DOCX {file_path}: {str(e)}")
            raise e

        return blocks

    @staticmethod
    def _format_excel_cell_value(val: Any) -> str:
        if val is None:
            return ""
        if isinstance(val, (datetime.datetime, datetime.date)):
            if isinstance(val, datetime.datetime) and (val.hour != 0 or val.minute != 0 or val.second != 0):
                return val.strftime("%Y-%m-%d %H:%M:%S")
            return val.strftime("%Y-%m-%d")
            
        s_val = str(val).strip()
        if not s_val:
            return ""
            
        # Trim trailing 00:00:00 from string datetimes
        s_val = re.sub(r'(\d{4}-\d{2}-\d{2})\s+00:00:00', r'\1', s_val)
        
        # Round floating point noise (e.g. 0.9999999999999999 -> 1, 1.4999999999999998 -> 1.5)
        if isinstance(val, float) or re.match(r'^-?\d+\.\d+$', s_val):
            try:
                f_val = float(s_val)
                rounded = round(f_val, 4)
                if rounded == int(rounded):
                    s_val = str(int(rounded))
                else:
                    s_val = f"{rounded:g}"
            except Exception:
                pass
                
        return s_val.replace('\n', ' ').replace('|', '\\|')

    @staticmethod
    def _is_section_title(row: tuple) -> bool:
        """
        Strict Section Title matcher.
        Only returns True if row is a genuine section title / banner.
        Does NOT match bottom notes, comments, or footnotes.
        """
        if not row:
            return False
        non_empty = [(i, str(c).strip()) for i, c in enumerate(row) if c is not None and str(c).strip()]
        if not non_empty or len(non_empty) > 3:
            return False
        
        first_idx, first_val = non_empty[0]
        lower_val = first_val.lower()

        # Bottom notes / comments start with explicit keywords or lowercase notes
        if any(lower_val.startswith(k) for k in ["ghi chú", "chú thích", "note", "chữ xanh", "*", "seminar", "mai show", "hướng dẫn"]):
            return False

        # If a row has multiple non-empty cells (e.g. ['Tổng số báo cáo', '690']), it is a Data/Key-Value row, NOT a Section Title banner!
        if len(non_empty) > 1:
            return False

        # Single cell section title rules:
        # 1. Decorative banners e.g. "━━ Phân loại... ━━", "📊 Thống kê..."
        if any(k in first_val for k in ["━━", "══", "──", "---"]) or first_val.startswith(("📊", "📋", "📌", "📁", "📑", "💡", "🔍")):
            return True

        # 2. Numbered section title: "1. PHÂN BỔ...", "2. WORKLOAD...", "Section 1", "Table 2", "Checklist..."
        if re.match(r'^(?:\d+[\.\:]|Phần|Bảng|Section|Table|Checklist)\s*', first_val, re.IGNORECASE):
            return True
            
        # 3. ALL CAPS banner text like "CHECKLIST TỪNG DỰ ÁN – NGUỒN GỐC"
        if first_val.isupper() and len(first_val) >= 5:
            return True
            
        # 4. Short concise single-cell title header (< 60 chars) with title case (no period, no equation)
        if len(first_val) < 60 and not first_val.endswith('.') and not re.search(r'\=\s*\d+', first_val):
            return True

        return False

    @staticmethod
    def _is_header_row_candidate(row: tuple) -> Tuple[bool, int]:
        if not row:
            return False, -999
        non_empty = [str(c).strip() for c in row if c is not None and str(c).strip()]
        if len(non_empty) < 2:
            return False, -999
        title_like = 0
        pure_num = 0
        for val in non_empty:
            if val.lower().startswith(("http://", "https://", "www.")):
                continue
            if len(val) < 50 and not re.match(r'^\d+(\.\d+)?%?$', val):
                title_like += 1
            elif re.match(r'^\d+(\.\d+)?%?$', val):
                pure_num += 1
        score = title_like * 2 - pure_num * 3
        is_candidate = title_like >= 2 and pure_num <= title_like
        return is_candidate, score

    @classmethod
    def _split_sheet_into_tables(cls, rows_list: List[tuple]) -> List[Dict[str, Any]]:
        tables = []
        current_title = ""
        current_rows = []
        
        for r_idx, row in enumerate(rows_list):
            if not row or not any(c is not None and str(c).strip() for c in row):
                if current_rows:
                    tables.append({
                        "title": current_title,
                        "rows": current_rows
                    })
                    current_rows = []
                    current_title = ""
                continue
            
            if cls._is_section_title(row):
                if current_rows:
                    tables.append({
                        "title": current_title,
                        "rows": current_rows
                    })
                    current_rows = []
                non_empty_vals = [str(c).strip() for c in row if c is not None and str(c).strip()]
                current_title = " - ".join(non_empty_vals)
                continue
            
            if current_rows and len(current_rows) >= 2:
                is_cand, score = cls._is_header_row_candidate(row)
                prev_cand, prev_score = cls._is_header_row_candidate(current_rows[-1])
                if is_cand and not prev_cand and score > 3:
                    tables.append({
                        "title": current_title,
                        "rows": current_rows
                    })
                    current_rows = []
                    current_title = ""

            current_rows.append(row)
            
        if current_rows:
            tables.append({
                "title": current_title,
                "rows": current_rows
            })
            
        return tables

    @staticmethod
    def parse_xlsx(file_path: str) -> List[NormalizedBlock]:
        logger.info(f"Parsing XLSX file: {file_path}")
        blocks: List[NormalizedBlock] = []
        doc_basename = os.path.splitext(os.path.basename(file_path))[0]
        base_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "../../storage/extracted_images"))
        os.makedirs(base_dir, exist_ok=True)
        source_order = 0

        # Step 1: Pre-map embedded images to exact Sheet & Row anchors using openpyxl drawings
        sheet_images_by_row: Dict[str, Dict[int, List[Dict[str, Any]]]] = {}
        doc_image_cache: Dict[str, Dict[str, Any]] = {}
        try:
            wb_draw = openpyxl.load_workbook(file_path, data_only=True)
            for s_name in wb_draw.sheetnames:
                ws = wb_draw[s_name]
                images = getattr(ws, "_images", [])
                sheet_images_by_row[s_name] = {}
                for idx, img in enumerate(images, start=1):
                    anchor_cell = getattr(img.anchor, "_from", None)
                    r_idx = anchor_cell.row + 1 if anchor_cell else 1
                    ext = getattr(img, 'format', 'png').lower()
                    try:
                        img_bytes = img._data()
                        img_info = cls._save_image_with_hash(
                            img_bytes=img_bytes,
                            doc_basename=f"{doc_basename}_{s_name}",
                            ext=ext,
                            base_dir=base_dir,
                            doc_image_cache=doc_image_cache,
                            alt_text=f"Embedded Image {idx} in Sheet {s_name}"
                        )
                        if r_idx not in sheet_images_by_row[s_name]:
                            sheet_images_by_row[s_name][r_idx] = []
                        sheet_images_by_row[s_name][r_idx].append({
                            "rel_path": img_info["rel_path"],
                            "full_path": img_info["full_path"],
                            "alt": img_info["alt"]
                        })
                    except Exception as e_save:
                        logger.warning(f"Could not save XLSX image: {e_save}")
        except Exception as draw_err:
            logger.warning(f"Fallback: openpyxl drawing inspection failed for {file_path}: {draw_err}")

        try:
            wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                
                # Skip hidden sheets (sheets hidden by users in MS Excel UI)
                if hasattr(sheet, 'sheet_state') and sheet.sheet_state and sheet.sheet_state != 'visible':
                    logger.info(f"Skipping hidden sheet '{sheet_name}' in XLSX file: {file_path}")
                    continue

                curr_sheet_imgs = sheet_images_by_row.get(sheet_name, {})
                
                rows_list = list(sheet.iter_rows(values_only=True))
                if not rows_list:
                    continue

                headers = []
                best_header_row_idx = 0
                max_cols_count = 0
                
                # Scan top 15 rows to find the True Header row with concise field titles (avoiding URL/metadata rows)
                try:
                    best_score = -999
                    for r_idx, r in enumerate(rows_list[:15]):
                        if not r:
                            continue
                        non_empty_cols = [i for i, c in enumerate(r) if c is not None and str(c).strip()]
                        if not non_empty_cols:
                            continue

                        # Count concise field title cells (< 70 chars, not URLs)
                        field_title_cols = [
                            i for i in non_empty_cols 
                            if not str(r[i]).strip().lower().startswith(("http://", "https://", "www."))
                            and len(str(r[i]).strip()) < 70
                        ]
                        url_count = len(non_empty_cols) - len(field_title_cols)
                        score = len(field_title_cols) * 2 - url_count * 3

                        # Calculate max column index across top 100 rows in the sheet to prevent truncating offset tables
                        sheet_max_col = 0
                        for r_sample in rows_list[:100]:
                            if r_sample:
                                ne = [i_c for i_c, v_c in enumerate(r_sample) if v_c is not None and str(v_c).strip()]
                                if ne:
                                    sheet_max_col = max(sheet_max_col, max(ne) + 1)

                        if score > best_score or (score == best_score and len(non_empty_cols) > max_cols_count):
                            best_score = score
                            max_cols_count = len(non_empty_cols)
                            best_header_row_idx = r_idx
                            max_col_idx = min(max(sheet_max_col, max(non_empty_cols) + 1), 35) # Cap at 35 columns max

                            raw_headers = []
                            for i in range(max_col_idx):
                                if i < len(r) and r[i] is not None and str(r[i]).strip():
                                    val = str(r[i]).strip().replace('\n', ' ').replace('|', '\\|')
                                    if val.lower().startswith(("http://", "https://", "www.")):
                                        val = "[Link]"
                                    elif len(val) > 40:
                                        val = val[:37] + "..."
                                    raw_headers.append(val)
                                else:
                                    raw_headers.append("")
                            headers = raw_headers
                except Exception as h_err:
                    logger.warning(f"Error finding table header in sheet '{sheet_name}': {h_err}")
                    continue

                # Detect if this sheet is a Timeline / Action Plan sheet (>12 cols & date/week headers present)
                is_timeline_sheet = False
                month_headers = {}
                timeline_col_map = {}
                core_headers_set = set()

                if max_cols_count >= 12:
                    # Scan row 3 for month group headers
                    current_month = ""
                    row3 = rows_list[2] if len(rows_list) > 2 else []
                    for idx_c, val_c in enumerate(row3):
                        if val_c is not None and str(val_c).strip():
                            val_str = str(val_c).strip()
                            if re.search(r'T\d{1,2}/\d{4}', val_str, re.IGNORECASE):
                                current_month = val_str
                        month_headers[idx_c] = current_month

                    # Scan top 10 rows for date/week columns & core column keywords
                    date_col_count = 0
                    curr_base_date = None

                    for idx_c in range(4, max_cols_count):
                        # Check top 10 rows for datetime base objects
                        for r_sample in rows_list[:10]:
                            if idx_c < len(r_sample) and isinstance(r_sample[idx_c], (datetime.datetime, datetime.date)):
                                curr_base_date = r_sample[idx_c]
                                break

                        # Check if any row in top 10 has day number 1..31
                        day_num_str = None
                        for r_sample in rows_list[:10]:
                            if idx_c < len(r_sample) and r_sample[idx_c] is not None:
                                v_s = str(r_sample[idx_c]).strip()
                                if re.match(r'^\d{1,2}$', v_s) and 1 <= int(v_s) <= 31:
                                    day_num_str = v_s
                                    break

                        if curr_base_date and day_num_str:
                            try:
                                d_num = int(day_num_str)
                                resolved_date = datetime.date(curr_base_date.year, curr_base_date.month, d_num).strftime("%d/%m/%Y")
                                timeline_col_map[idx_c] = {"col_name": resolved_date, "month": ""}
                                date_col_count += 1
                            except Exception:
                                pass

                    for r_sample in rows_list[:10]:
                        if not r_sample:
                            continue
                        for idx_c, val_c in enumerate(r_sample):
                            if val_c is None or not str(val_c).strip():
                                continue
                            val_str = str(val_c).strip()

                            # Check core task keywords
                            if idx_c < 10 and any(kw in val_str.upper() for kw in ["NO", "STT", "PROJECT", "DỰ ÁN", "TASK", "CÔNG VIỆC", "TEAM", "PERSON", "PIC", "STATUS", "TRẠNG THÁI", "HẠNG MỤC"]):
                                core_headers_set.add(val_str)

                            is_date_col = False
                            date_formatted = None

                            if isinstance(val_c, (datetime.datetime, datetime.date)):
                                is_date_col = True
                                date_formatted = val_c.strftime("%d/%m/%Y")
                            elif re.search(r'\d{1,2}-[Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec]{3}-\d{2,4}', val_str, re.IGNORECASE) or re.search(r'\d{4}-\d{2}-\d{2}', val_str) or re.search(r'\d{1,2}/\d{1,2}/\d{2,4}', val_str) or re.match(r'W\d{1,2}', val_str, re.IGNORECASE):
                                is_date_col = True
                                date_formatted = val_str.split(" ")[0]

                            if is_date_col and idx_c >= 4:
                                date_col_count += 1
                                if idx_c not in timeline_col_map:
                                    m_str = month_headers.get(idx_c, "")
                                    timeline_col_map[idx_c] = {
                                        "col_name": date_formatted or val_str or f"W{idx_c+1}",
                                        "month": m_str
                                    }

                    if (date_col_count >= 4 or len(timeline_col_map) >= 4) and len(core_headers_set) >= 1:
                        is_timeline_sheet = True

                # ==================================================
                # MODE A: TIMELINE AGGREGATION MODE (For Task Plans)
                # ==================================================
                if is_timeline_sheet:
                    logger.info(f"Timeline Action Plan sheet detected in '{sheet_name}'. Activating Timeline Aggregation Engine.")
                    
                    # Dynamically map column positions by scanning ALL top 10 rows
                    col_map_idx = {}
                    
                    # Find first populated column index in sample data rows
                    first_col_idx = 0
                    for r_sample in rows_list[best_header_row_idx+1:best_header_row_idx+10]:
                        if not r_sample:
                            continue
                        non_empty_indices = [i_c for i_c, v_c in enumerate(r_sample) if v_c is not None and str(v_c).strip()]
                        if non_empty_indices:
                            first_col_idx = min(non_empty_indices)
                            break

                    for r_sample in rows_list[:10]:
                        if not r_sample:
                            continue
                        for c_idx, c_val in enumerate(r_sample):
                            if c_val is None:
                                continue
                            c_lower = str(c_val).strip().lower()
                            if any(k in c_lower for k in ["stt", "no.", "no"]) and c_idx < 5 and "stt" not in col_map_idx:
                                col_map_idx["stt"] = c_idx
                            elif c_idx >= 2 and any(k in c_lower for k in ["hạng mục", "công việc", "task title", "task", "phase", "description", "title", "project"]):
                                if "task_name" not in col_map_idx and c_idx != col_map_idx.get("stt"):
                                    col_map_idx["task_name"] = c_idx
                            elif any(k in c_lower for k in ["person in charge", "pic", "owner", "người làm", "phụ trách"]):
                                if "pic" not in col_map_idx:
                                    col_map_idx["pic"] = c_idx
                            elif any(k in c_lower for k in ["team in charge", "team"]):
                                if "team" not in col_map_idx:
                                    col_map_idx["team"] = c_idx
                            elif any(k in c_lower for k in ["ưu tiên", "priority"]):
                                if "priority" not in col_map_idx:
                                    col_map_idx["priority"] = c_idx
                            elif any(k in c_lower for k in ["trạng thái", "status"]):
                                if "status" not in col_map_idx:
                                    col_map_idx["status"] = c_idx
                            elif any(k in c_lower for k in ["tiến độ", "pct", "%", "complete", "hoàn thành"]):
                                if "progress" not in col_map_idx:
                                    col_map_idx["progress"] = c_idx

                    # Smart fallback for STT & Task Name if not explicitly matched by header keywords
                    if "stt" not in col_map_idx:
                        col_map_idx["stt"] = first_col_idx

                    if "task_name" not in col_map_idx:
                        pic_idx = col_map_idx.get("pic", -1)
                        stt_idx = col_map_idx.get("stt", -1)
                        possible_task_cols = [c for c in range(first_col_idx, min(first_col_idx + 6, max_cols_count)) if c != pic_idx and c != stt_idx]
                        
                        col_total_vol = {}
                        for c_cand in possible_task_cols:
                            lens = [len(str(r[c_cand]).strip()) for r in rows_list[best_header_row_idx+1:best_header_row_idx+35] if c_cand < len(r) and r[c_cand] is not None and str(r[c_cand]).strip()]
                            col_total_vol[c_cand] = sum(lens)
                            
                        best_cand = max(col_total_vol.items(), key=lambda x: x[1])[0] if col_total_vol else first_col_idx + 1
                        col_map_idx["task_name"] = best_cand

                    # Default fallback indexes
                    idx_stt = col_map_idx["stt"]
                    idx_task = col_map_idx["task_name"]
                    idx_pic = col_map_idx.get("pic", None)
                    idx_team = col_map_idx.get("team", None)
                    idx_prio = col_map_idx.get("priority", None)
                    idx_stat = col_map_idx.get("status", None)
                    idx_prog = col_map_idx.get("progress", None)

                    # Also inspect cell fills from openpyxl sheet object for colored Gantt bars
                    try:
                        ws_openpyxl = wb_draw[sheet_name]
                        ws_openpyxl_rows = list(ws_openpyxl.iter_rows())
                    except Exception:
                        ws_openpyxl_rows = []
                    # Calculate column baseline fill colors to distinguish Gantt bars from vertical column themes/stripes
                    col_default_fills = {}
                    if ws_openpyxl_rows and timeline_col_map:
                        for c_k in timeline_col_map.keys():
                            c_fills = []
                            for r_o in ws_openpyxl_rows[best_header_row_idx + 1:min(len(ws_openpyxl_rows), best_header_row_idx + 50)]:
                                if c_k < len(r_o):
                                    cell_o = r_o[c_k]
                                    if cell_o.fill and cell_o.fill.fill_type and cell_o.fill.start_color:
                                        sc = cell_o.fill.start_color
                                        f_val = str(sc.rgb) if sc.type == "rgb" else (f"THEME_{sc.theme}" if sc.type == "theme" else str(sc.value))
                                        c_fills.append(f_val)
                                    else:
                                        c_fills.append("NONE")
                            most_c = Counter(c_fills).most_common(1)[0][0] if c_fills else "NONE"
                            col_default_fills[c_k] = most_c

                    timeline_records = []
                    raw_row_index = best_header_row_idx + 1

                    for row_i, row in enumerate(rows_list[best_header_row_idx + 1:], start=best_header_row_idx + 1):
                        raw_row_index += 1
                        if not row:
                            continue

                        stt = str(row[idx_stt]).strip() if idx_stt is not None and idx_stt < len(row) and row[idx_stt] is not None else ""
                        task_name = str(row[idx_task]).strip() if idx_task is not None and idx_task < len(row) and row[idx_task] is not None else ""
                        if not task_name and idx_task is not None and (idx_task + 1) < len(row) and row[idx_task + 1] is not None:
                            task_name = str(row[idx_task + 1]).strip()
                        pic_val = str(row[idx_pic]).strip() if idx_pic is not None and idx_pic < len(row) and row[idx_pic] is not None else ""
                        team_val = str(row[idx_team]).strip() if idx_team is not None and idx_team < len(row) and row[idx_team] is not None else ""
                        
                        if pic_val and team_val and pic_val != team_val:
                            pic = f"{pic_val} ({team_val})"
                        else:
                            pic = pic_val or team_val or "Unassigned"

                        priority = str(row[idx_prio]).strip() if idx_prio is not None and idx_prio < len(row) and row[idx_prio] is not None else "Normal"
                        status_raw = str(row[idx_stat]).strip() if idx_stat is not None and idx_stat < len(row) and row[idx_stat] is not None else ""
                        progress_raw = str(row[idx_prog]).strip() if idx_prog is not None and idx_prog < len(row) and row[idx_prog] is not None else ""

                        if not task_name or task_name.lower() in ["stt", "no.", "task title", "hạng mục / công việc"]:
                            continue

                        progress_val = None
                        if progress_raw:
                            try:
                                p_num = float(progress_raw.replace('%', '').strip())
                                progress_val = p_num if p_num <= 1.0 else p_num / 100.0
                            except ValueError:
                                pass

                        if progress_val is None and status_raw:
                            try:
                                p_num = float(status_raw.replace('%', '').strip())
                                progress_val = p_num if p_num <= 1.0 else p_num / 100.0
                            except ValueError:
                                pass

                        if status_raw and not status_raw.replace('.', '', 1).replace('%', '').strip().isdigit():
                            s_clean = status_raw.lower()
                            if any(k in s_clean for k in ["done", "completed", "hoàn thành"]):
                                status = "Completed"
                            elif any(k in s_clean for k in ["inprocess", "in progress", "in-progress", "đang làm"]):
                                status = "In Progress"
                            elif any(k in s_clean for k in ["hold", "tạm dừng", "delay"]):
                                status = "On Hold"
                            elif any(k in s_clean for k in ["pending", "chờ", "chưa"]):
                                status = "Pending"
                            else:
                                status = status_raw
                        elif progress_val is not None:
                            if progress_val >= 1.0:
                                status = "Completed"
                            elif progress_val > 0.0:
                                status = "In Progress"
                            else:
                                status = "Pending"
                        else:
                            status = "Pending"

                        if progress_val is not None:
                            progress_str = f"{int(round(progress_val * 100))}%"
                        else:
                            progress_str = "0%"

                        # Inspect cell values AND cell fill colors for active timeline markers
                        active_cols = []
                        cell_objects = ws_openpyxl_rows[row_i] if row_i < len(ws_openpyxl_rows) else []

                        # Check if current row is a Category/Header row (uniform fill across >75% timeline columns)
                        ignored_bg_fills = ["NONE", "00000000", "FFFFFFFF", "00FFFFFF", "FFCCCCCC", "FFF2F2F2", "FFD0D0D0", "FFE0E0E0", "FFF9F9F9", "FFEEEEEE", "FFFAFAFA"]
                        row_t_fills = []
                        for c_k in timeline_col_map.keys():
                            if c_k < len(cell_objects):
                                cell_o = cell_objects[c_k]
                                if cell_o.fill and cell_o.fill.fill_type and cell_o.fill.start_color:
                                    sc = cell_o.fill.start_color
                                    f_val = str(sc.rgb) if sc.type == "rgb" else (f"THEME_{sc.theme}" if sc.type == "theme" else str(sc.value))
                                    row_t_fills.append(f_val)
                                else:
                                    row_t_fills.append("NONE")

                        freq_fill, freq_cnt = Counter(row_t_fills).most_common(1)[0] if row_t_fills else ("NONE", 0)
                        is_section_header = (
                            (freq_cnt / len(row_t_fills) > 0.75) 
                            and freq_fill not in ignored_bg_fills
                            and any(freq_fill != col_default_fills.get(c_k) for c_k in timeline_col_map.keys())
                        ) if row_t_fills else False

                        for c_idx, t_info in timeline_col_map.items():
                            is_active = False
                            if c_idx < len(row) and row[c_idx] is not None and str(row[c_idx]).strip():
                                c_val = str(row[c_idx]).strip().lower()
                                if c_val in ["x", "1", "true", "yes", "done", "100%", "in progress"]:
                                    is_active = True
                            
                            # Check cell fill color if value check didn't trigger AND row is not a full section header fill
                            if not is_active and not is_section_header and c_idx < len(cell_objects):
                                cell_obj = cell_objects[c_idx]
                                if cell_obj.fill and cell_obj.fill.fill_type and cell_obj.fill.start_color:
                                    sc = cell_obj.fill.start_color
                                    rgb = str(sc.rgb) if sc.type == "rgb" else (f"THEME_{sc.theme}" if sc.type == "theme" else str(sc.value))
                                    # Ignore standard empty background fills
                                    if rgb not in ignored_bg_fills:
                                        # Must be different from the column's default baseline fill!
                                        if rgb != col_default_fills.get(c_idx):
                                            is_active = True

                            if is_active:
                                active_cols.append((c_idx, t_info))

                        if active_cols:
                            start_date = active_cols[0][1]["col_name"]
                            end_date = active_cols[-1][1]["col_name"]
                            start_month = active_cols[0][1].get("month", "")
                            end_month = active_cols[-1][1].get("month", "")

                            if start_month and end_month:
                                month_range_str = f" ({start_month})" if start_month == end_month else f" ({start_month} – {end_month})"
                            elif start_month or end_month:
                                month_range_str = f" ({start_month or end_month})"
                            else:
                                month_range_str = ""

                            timeline_summary = f"{start_date} – {end_date}{month_range_str}".strip() if start_date != end_date else f"{start_date}{month_range_str}".strip()
                        else:
                            timeline_summary = "Chưa xếp lịch"

                        rec_md = f"- **[{stt}] {task_name}**\n" \
                                 f"  • **PIC:** {pic} | **Ưu tiên:** {priority} | **Trạng thái:** {status} | **Tiến độ:** {progress_str}\n" \
                                 f"  • **Thời gian thực hiện:** {timeline_summary}"
                        timeline_records.append(rec_md)

                    # Group timeline task records into chunks
                    t_buffer = []
                    t_len = 0
                    for rec in timeline_records:
                        if t_len + len(rec) > 1150 and t_buffer:
                            source_order += 1
                            blocks.append(NormalizedBlock(
                                block_id=str(uuid.uuid4()),
                                source_type="xlsx",
                                block_type="paragraph",
                                content=f"### {sheet_name} (Kế hoạch hành động)\n\n" + "\n\n".join(t_buffer),
                                heading_path=[f"Sheet: {sheet_name}"],
                                source_order=source_order
                            ))
                            t_buffer = []
                            t_len = 0

                        t_buffer.append(rec)
                        t_len += len(rec)

                    if t_buffer:
                        source_order += 1
                        blocks.append(NormalizedBlock(
                            block_id=str(uuid.uuid4()),
                            source_type="xlsx",
                            block_type="paragraph",
                            content=f"### {sheet_name} (Kế hoạch hành động)\n\n" + "\n\n".join(t_buffer),
                            heading_path=[f"Sheet: {sheet_name}"],
                            source_order=source_order
                        ))
                    continue

                # Check if this sheet is a Large Data Dump Sheet (> 200 rows OR > 10,000 total characters)
                total_sheet_rows = len(rows_list)
                total_sheet_chars = sum(len(str(c)) for r in rows_list for c in r if c is not None)
                is_large_sheet = total_sheet_rows > 200 or total_sheet_chars > 10000

                if is_large_sheet:
                    logger.warning(
                        f"Sheet '{sheet_name}' in {os.path.basename(file_path)} exceeds 200 rows / 10K chars "
                        f"({total_sheet_rows} rows, {total_sheet_chars} chars). Filtering to keep 1ST SINGLE CHUNK ONLY."
                    )

                # ==================================================
                # MODE B: STANDARD MARKDOWN PIPE TABLE MODE (MULTI-TABLE SEGMENTED)
                # ==================================================
                table_segments = ParserService._split_sheet_into_tables(rows_list)
                logger.info(f"Segmented sheet '{sheet_name}' into {len(table_segments)} table(s).")
                sheet_blocks = []

                for tbl_idx, tbl in enumerate(table_segments, start=1):
                    tbl_title = tbl.get("title", "")
                    tbl_rows = tbl.get("rows", [])
                    if not tbl_rows:
                        continue

                    # Find best header row inside this segment
                    best_header_idx = 0
                    best_score = -999
                    for i, r in enumerate(tbl_rows[:5]):
                        is_cand, score = ParserService._is_header_row_candidate(r)
                        if score > best_score:
                            best_score = score
                            best_header_idx = i

                    header_row = tbl_rows[best_header_idx]
                    data_rows = tbl_rows[best_header_idx + 1:]

                    if is_large_sheet:
                        # Cap data_rows to top 15 rows so entire table + note fits inside EXACTLY 1 single NormalizedBlock
                        data_rows = data_rows[:15]

                    # Active column indices for this segment
                    max_c = max(len(r) for r in tbl_rows)
                    active_col_indices = []
                    for c_i in range(max_c):
                        has_val = any(
                            c_i < len(r) and r[c_i] is not None and str(r[c_i]).strip() != ""
                            for r in tbl_rows
                        )
                        if has_val:
                            active_col_indices.append(c_i)

                    if not active_col_indices:
                        active_col_indices = list(range(min(max_c, 5)))

                    if is_large_sheet and len(active_col_indices) > 12:
                        # Cap active columns to top 12 columns max for large sheet representative sample to prevent block explosion
                        active_col_indices = active_col_indices[:12]

                    # Detect if segment is a 2-column Key-Value / Summary List vs a True Grid Table
                    is_key_value = len(active_col_indices) <= 2 and (best_score <= 2 or len(active_col_indices) == 1)

                    if is_key_value:
                        # Format as simple raw pipe-separated text lines (e.g. "a | b")
                        kv_lines = []
                        all_rows_to_process = tbl_rows if best_score <= 0 else tbl_rows[best_header_idx:]

                        for r in all_rows_to_process:
                            vals = [
                                ParserService._format_excel_cell_value(r[c_i])
                                if c_i < len(r) and r[c_i] is not None else ""
                                for c_i in active_col_indices
                            ]
                            non_empty_vals = [v for v in vals if v]
                            if not non_empty_vals:
                                continue
                            kv_lines.append(" | ".join(non_empty_vals))

                        if kv_lines:
                            title_prefix = f"### Sheet: {sheet_name} - {tbl_title}\n\n" if tbl_title else f"### Sheet: {sheet_name}\n\n"
                            kv_content = title_prefix + "\n".join(kv_lines)
                            source_order += 1
                            sheet_blocks.append(NormalizedBlock(
                                block_id=str(uuid.uuid4()),
                                source_type="xlsx",
                                block_type="paragraph",
                                content=kv_content.strip(),
                                heading_path=[f"Sheet: {sheet_name}", tbl_title or f"Tóm tắt {tbl_idx}"],
                                sheet_name=sheet_name,
                                table_id=f"sheet_{sheet_name}_kv_{tbl_idx}",
                                requires_llm_summary=False,
                                source_order=source_order
                            ))
                        continue

                    # Standard Multi-Column Pipe Table Mode
                    # Headers
                    segment_headers = []
                    for c_i in active_col_indices:
                        val = ""
                        if c_i < len(header_row) and header_row[c_i] is not None:
                            val = ParserService._format_excel_cell_value(header_row[c_i])
                        segment_headers.append(val if val else f"Col_{c_i+1}")

                    header_line = "| " + " | ".join(segment_headers) + " |\n"
                    divider_line = "| " + " | ".join(["---"] * len(active_col_indices)) + " |\n"
                    
                    if tbl_title:
                        title_prefix = f"### Sheet: {sheet_name} - {tbl_title}\n\n"
                    else:
                        title_prefix = f"### Sheet: {sheet_name} (Bảng {tbl_idx})\n\n"

                    header_md = title_prefix + header_line + divider_line

                    batch_rows = []
                    current_char_count = len(header_md)
                    data_rows_count = 0

                    for r in data_rows:
                        row_vals = []
                        for i in active_col_indices:
                            v_str = ParserService._format_excel_cell_value(r[i]) if i < len(r) and r[i] is not None else ""
                            if is_large_sheet and len(v_str) > 70:
                                v_str = v_str[:67] + "..."
                            row_vals.append(v_str)

                        if not any(row_vals):
                            continue

                        row_line = "| " + " | ".join(row_vals) + " |\n"
                        row_len = len(row_line)
                        data_rows_count += 1

                        batch_rows.append(row_line)
                        current_char_count += row_len

                        if is_large_sheet and (current_char_count >= 650 or len(batch_rows) >= 2):
                            # LARGE SHEET GUARDRAIL: Stop accumulating rows once character budget reaches ~650 chars so total block + note never exceeds 1200 chars
                            break

                        if not is_large_sheet and (current_char_count >= 1150 or len(batch_rows) >= 20):
                            source_order += 1
                            table_content = header_md + "".join(batch_rows)
                            sheet_blocks.append(NormalizedBlock(
                                block_id=str(uuid.uuid4()),
                                source_type="xlsx",
                                block_type="table",
                                content=table_content.strip(),
                                heading_path=[f"Sheet: {sheet_name}", tbl_title or f"Bảng {tbl_idx}"],
                                sheet_name=sheet_name,
                                table_id=f"sheet_{sheet_name}_tbl_{tbl_idx}",
                                requires_llm_summary=False,
                                source_order=source_order
                            ))
                            batch_rows = []
                            current_char_count = len(header_md)

                    if batch_rows:
                        source_order += 1
                        table_content = header_md + "".join(batch_rows)
                        if is_large_sheet:
                            remaining_rows = max(0, total_sheet_rows - data_rows_count)
                            table_content += f"\n\n*(Lưu ý: Sheet '{sheet_name}' có tổng cộng {total_sheet_rows:,} dòng ({total_sheet_chars:,} ký tự). Đã trích xuất {data_rows_count} dòng đại diện trong chunk này, còn {remaining_rows:,} dòng dữ liệu khác chưa hiển thị để tránh quá tải CSDL.)*"

                        sheet_blocks.append(NormalizedBlock(
                            block_id=str(uuid.uuid4()),
                            source_type="xlsx",
                            block_type="table",
                            content=table_content.strip(),
                            heading_path=[f"Sheet: {sheet_name}", tbl_title or f"Bảng {tbl_idx}"],
                            sheet_name=sheet_name,
                            table_id=f"sheet_{sheet_name}_tbl_{tbl_idx}",
                            requires_llm_summary=False,
                            source_order=source_order
                        ))

                    if is_large_sheet:
                        # Large sheet guardrail: Stop processing further tables in this sheet
                        break

                blocks.extend(sheet_blocks)

        except Exception as e:
            logger.error(f"Error parsing XLSX {file_path}: {str(e)}")
            raise e

        return blocks

    @classmethod
    def parse_pptx(cls, file_path: str) -> List[NormalizedBlock]:
        logger.info(f"Parsing PPTX file with Tables and Images: {file_path}")
        doc_basename = os.path.splitext(os.path.basename(file_path))[0]
        base_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "../../storage/extracted_images"))
        os.makedirs(base_dir, exist_ok=True)

        blocks: List[NormalizedBlock] = []
        doc_image_cache: Dict[str, Dict[str, Any]] = {}
        source_order = 0

        try:
            import pptx
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            prs = pptx.Presentation(file_path)
            slides_list = list(prs.slides)

            def extract_shapes(shapes_list, slide_num, heading_path, img_counter, slide_title_text):
                nonlocal source_order
                for shape in shapes_list:
                    # 1. Handle Group Shapes recursively
                    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        img_counter = extract_shapes(shape.shapes, slide_num, heading_path, img_counter, slide_title_text)
                        continue

                    # 2. Handle Tables in PPTX
                    if shape.has_table:
                        table = shape.table
                        table_rows = []
                        for row in table.rows:
                            row_vals = [cell.text.strip().replace('\n', ' ').replace('|', '\\|') for cell in row.cells]
                            table_rows.append(row_vals)
                        
                        if table_rows:
                            max_c = max(len(r) for r in table_rows)
                            active_cols = [c for c in range(max_c) if any(c < len(r) and r[c] != "" for r in table_rows)]
                            if not active_cols:
                                active_cols = list(range(max_c))
                            
                            headers = [table_rows[0][c] if c < len(table_rows[0]) else "" for c in active_cols]
                            header_line = "| " + " | ".join(headers) + " |\n"
                            divider_line = "| " + " | ".join(["---"] * len(active_cols)) + " |\n"
                            
                            body_lines = []
                            for r_vals in table_rows[1:]:
                                line_vals = [r_vals[c] if c < len(r_vals) else "" for c in active_cols]
                                if any(v != "" for v in line_vals):
                                    body_lines.append("| " + " | ".join(line_vals) + " |\n")
                            
                            table_md = header_line + divider_line + "".join(body_lines)
                            if table_md.strip():
                                source_order += 1
                                blocks.append(NormalizedBlock(
                                    block_id=str(uuid.uuid4()),
                                    source_type="pptx",
                                    block_type="table",
                                    content=table_md.strip(),
                                    heading_path=heading_path,
                                    page_start=slide_num,
                                    page_end=slide_num,
                                    requires_llm_summary=False,
                                    source_order=source_order
                                ))
                        continue

                    # 3. Handle Images in PPTX with Hash Deduplication
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or hasattr(shape, "image"):
                        try:
                            img = shape.image
                            ext = getattr(img, "ext", "png")
                            img_counter += 1
                            img_info = cls._save_image_with_hash(
                                img_bytes=img.blob,
                                doc_basename=doc_basename,
                                ext=ext,
                                base_dir=base_dir,
                                doc_image_cache=doc_image_cache,
                                alt_text=f"Embedded Image {img_counter} in Slide {slide_num}"
                            )

                            # If an image repeats > 5 times (decorative template logo), mark summary as False
                            is_frequent_logo = img_info.get("count", 1) > 5

                            img_tag = f"<img src='{img_info['rel_path']}' alt='{img_info['alt']}' />"
                            source_order += 1
                            blocks.append(NormalizedBlock(
                                block_id=str(uuid.uuid4()),
                                source_type="pptx",
                                block_type="image",
                                content=img_tag,
                                image_path=img_info["full_path"],
                                requires_llm_summary=not is_frequent_logo,
                                heading_path=heading_path,
                                page_start=slide_num,
                                page_end=slide_num,
                                source_order=source_order
                            ))
                        except Exception as e_img:
                            logger.warning(f"Could not extract PPTX image on slide {slide_num}: {e_img}")
                        continue

                    # 4. Handle Text Frames
                    if shape.has_text_frame:
                        txt = shape.text.strip()
                        if txt and txt != slide_title_text:
                            source_order += 1
                            blocks.append(NormalizedBlock(
                                block_id=str(uuid.uuid4()),
                                source_type="pptx",
                                block_type="paragraph",
                                content=txt,
                                heading_path=heading_path,
                                page_start=slide_num,
                                page_end=slide_num,
                                source_order=source_order
                            ))

                return img_counter

            for slide_idx, slide in enumerate(slides_list):
                slide_num = slide_idx + 1
                slide_title_text = ""
                try:
                    if slide.shapes.title and slide.shapes.title.text:
                        slide_title_text = slide.shapes.title.text.strip()
                except Exception:
                    pass

                heading_title = f"Slide {slide_num}: {slide_title_text}" if slide_title_text else f"Slide {slide_num}"
                heading_path = [heading_title]

                extract_shapes(slide.shapes, slide_num, heading_path, 0, slide_title_text)

        except Exception as e:
            logger.error(f"Error parsing PPTX {file_path}: {str(e)}")
            raise e

        logger.info(f"Parsed PPTX file {os.path.basename(file_path)}. Extracted {len(blocks)} blocks.")
        return blocks

    @classmethod
    def parse(cls, file_path: str, file_type: str) -> List[NormalizedBlock]:
        file_type = file_type.lower().strip(".")
        if file_type == "pdf":
            return cls.parse_pdf(file_path)
        elif file_type in ["docx", "doc"]:
            return cls.parse_docx(file_path)
        elif file_type in ["xlsx", "xls", "csv"]:
            return cls.parse_xlsx(file_path)
        elif file_type in ["pptx", "ppt"]:
            return cls.parse_pptx(file_path)
        elif file_type in ["md", "markdown"]:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
            return cls.parse_markdown_content(content, source_type="md")
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
